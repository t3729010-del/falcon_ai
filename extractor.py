from PyPDF2 import PdfReader

try:
    from pdf2image import convert_from_path
except ImportError:
    convert_from_path = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from PIL import Image
try:
    import pytesseract
except ImportError:
    pytesseract = None

def extract_pdf_text(pdf_path):

    text = ""

    try:

        print("PDF PATH:",pdf_path)

        reader = PdfReader(pdf_path)

        for page in reader.pages:

            page_text = page.extract_text()

            if page_text:
                text += page_text

    except Exception as e:
        print("PDF Extraction Error:",e)

    if text.strip():
        return text

    try:

        print("Trying OCR...")

        if convert_from_path is None:
            raise Exception("pdf2image not installed")

        images = convert_from_path(pdf_path)

        print("Images Created:", len(images))

        for i, image in enumerate(images):

            page_text = pytesseract.image_to_string(
                image,
                lang="eng"
            )   

            print(
                f"Page {i+1} OCR Length:",
                len(page_text)
            )

            text += page_text

        print(
            "FINAL OCR LENGTH:",
            len(text)
        )

    except Exception as e:

        print("OCR ERROR:")
        print(str(e))

    return(text)   

def extract_docx_text(path):

    if Document is None:
        return ""

    doc = Document(path)

    text = ""

    for paragraph in doc.paragraphs:

        text += paragraph.text + "\n"

    return text

def extract_pptx_text(path):

    if Presentation is None:
        return ""

    prs = Presentation(path)

    text = ""

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text += shape.text + "\n"

    return text

def extract_image_text(path):

    if pytesseract is None:
        return ""

    image = Image.open(path)

    return pytesseract.image_to_string(
        image,
        lang="eng+urd+ara"
    )

def extract_text(file_path, extension):

    extension = extension.lower()

    if extension == "pdf":
        return extract_pdf_text(file_path)

    elif extension == "docx":
        return extract_docx_text(file_path)

    elif extension in ["ppt","pptx"]:
        return extract_pptx_text(file_path)

    elif extension in [
        "jpg",
        "jpeg",
        "png"
    ]:
        return extract_image_text(file_path)
    
    elif extension == "txt":

        with open(
            file_path,
            "r",
            encoding="utf-8",
            errors="ignore"
        ) as f:

            return f.read()

    return ""